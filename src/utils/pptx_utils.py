import io
from pptx import Presentation
from sb_utils.logger_utils import logger

def extract_text_from_pptx(file_stream: io.BytesIO) -> str:
    """
    Extracts text from a .pptx file stream.
    """
    try:
        presentation = Presentation(file_stream)
        full_text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        
        logger.info("Successfully extracted text from PPTX file.")
        return '\\n'.join(full_text)
    except Exception as e:
        logger.error(f"Could not read PPTX file: {e}", exc_info=True)
        raise ValueError("Invalid or corrupted PPTX file.")
