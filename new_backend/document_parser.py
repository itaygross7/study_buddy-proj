# document_parsers.py
"""
Document parsers.

Responsibilities:
- Take a FileRef and return a normalized DocumentText.
- Support: TXT, MD, PDF, DOCX, PPTX, HTML, EPUB, images (OCR).
"""

from __future__ import annotations

import logging
import pickle
from abc import ABC, abstractmethod
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List

from file_utils import FileRef, FileType
from models_utils import InfraError

logger = logging.getLogger("studybuddy.parsers")


@dataclass
class DocumentText:
    """
    Normalized representation of a document after parsing.
    """
    text: str
    chunks: List[str]
    metadata: Dict[str, str]


# ---------- helper: simple chunker ----------

def simple_text_chunker(text: str, max_chars: int = 4_000) -> List[str]:
    if not text:
        return []

    paragraphs = text.split("\n\n")
    chunks: List[str] = []
    current: List[str] = []
    current_len = 0

    for para in paragraphs:
        p = para.strip()
        if not p:
            continue
        if current_len + len(p) + 2 > max_chars and current:
            chunks.append("\n\n".join(current))
            current = [p]
            current_len = len(p)
        else:
            current.append(p)
            current_len += len(p) + 2

    if current:
        chunks.append("\n\n".join(current))

    return chunks


# ---------- base parser ----------

class BaseDocumentParser(ABC):
    file_types: List[FileType] = []

    @classmethod
    def can_handle(cls, file_ref: FileRef) -> bool:
        return file_ref.file_type in cls.file_types

    @abstractmethod
    def parse(self, file_ref: FileRef) -> DocumentText:  # pragma: no cover
        ...


# ---------- plain text ----------

class PlainTextParser(BaseDocumentParser):
    file_types = [FileType.TXT]

    def parse(self, file_ref: FileRef) -> DocumentText:
        text = file_ref.path.read_text(encoding="utf-8", errors="ignore")
        chunks = simple_text_chunker(text)
        meta = {
            "file_id": file_ref.id,
            "file_type": file_ref.file_type.value,
            "parser": "PlainTextParser",
        }
        return DocumentText(text=text, chunks=chunks, metadata=meta)


# ---------- markdown ----------

class MarkdownParser(BaseDocumentParser):
    file_types = [FileType.MD]

    def parse(self, file_ref: FileRef) -> DocumentText:
        text = file_ref.path.read_text(encoding="utf-8", errors="ignore")
        chunks = simple_text_chunker(text)
        meta = {
            "file_id": file_ref.id,
            "file_type": file_ref.file_type.value,
            "parser": "MarkdownParser",
        }
        return DocumentText(text=text, chunks=chunks, metadata=meta)


# ---------- PDF ----------

class PDFParser(BaseDocumentParser):
    file_types = [FileType.PDF]

    def parse(self, file_ref: FileRef) -> DocumentText:
        try:
            import PyPDF2  # type: ignore
        except ImportError as exc:  # pragma: no cover
            raise InfraError("PyPDF2 is required to parse PDF files.") from exc

        text_parts: List[str] = []
        with file_ref.path.open("rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                try:
                    text_parts.append(page.extract_text() or "")
                except Exception:
                    continue

        full_text = "\n\n".join(t for t in text_parts if t)
        chunks = simple_text_chunker(full_text)
        meta = {
            "file_id": file_ref.id,
            "file_type": file_ref.file_type.value,
            "parser": "PDFParser",
            "pages": str(len(text_parts)),
        }
        return DocumentText(text=full_text, chunks=chunks, metadata=meta)


# ---------- DOCX (Word) ----------

class DOCXParser(BaseDocumentParser):
    file_types = [FileType.DOCX]

    def parse(self, file_ref: FileRef) -> DocumentText:
        try:
            import docx  # type: ignore
        except ImportError as exc:  # pragma: no cover
            raise InfraError("python-docx is required to parse DOCX files.") from exc

        doc = docx.Document(str(file_ref.path))
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        full_text = "\n\n".join(paragraphs)
        chunks = simple_text_chunker(full_text)

        meta = {
            "file_id": file_ref.id,
            "file_type": file_ref.file_type.value,
            "parser": "DOCXParser",
            "paragraphs": str(len(paragraphs)),
        }
        return DocumentText(text=full_text, chunks=chunks, metadata=meta)


# ---------- PPTX (PowerPoint) ----------

class PPTXParser(BaseDocumentParser):
    file_types = [FileType.PPTX]

    def parse(self, file_ref: FileRef) -> DocumentText:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as exc:  # pragma: no cover
            raise InfraError("python-pptx is required to parse PPTX files.") from exc

        prs = Presentation(str(file_ref.path))
        slide_texts: List[str] = []
        for slide in prs.slides:
            parts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
            if parts:
                slide_texts.append("\n".join(parts))

        full_text = "\n\n--- SLIDE BREAK ---\n\n".join(slide_texts)
        chunks = simple_text_chunker(full_text)
        meta = {
            "file_id": file_ref.id,
            "file_type": file_ref.file_type.value,
            "parser": "PPTXParser",
            "slides": str(len(slide_texts)),
        }
        return DocumentText(text=full_text, chunks=chunks, metadata=meta)


# ---------- HTML ----------

class HTMLParser(BaseDocumentParser):
    file_types = [FileType.HTML]

    def parse(self, file_ref: FileRef) -> DocumentText:
        try:
            from bs4 import BeautifulSoup  # type: ignore
        except ImportError as exc:  # pragma: no cover
            raise InfraError("beautifulsoup4 is required to parse HTML files.") from exc

        raw_html = file_ref.path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(raw_html, "html.parser")

        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()

        text = soup.get_text(separator="\n")
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        full_text = "\n".join(lines)
        chunks = simple_text_chunker(full_text)

        meta = {
            "file_id": file_ref.id,
            "file_type": file_ref.file_type.value,
            "parser": "HTMLParser",
        }
        return DocumentText(text=full_text, chunks=chunks, metadata=meta)


# ---------- EPUB ----------

class EPUBParser(BaseDocumentParser):
    file_types = [FileType.EPUB]

    def parse(self, file_ref: FileRef) -> DocumentText:
        try:
            from ebooklib import epub  # type: ignore
            from bs4 import BeautifulSoup  # type: ignore
        except ImportError as exc:  # pragma: no cover
            raise InfraError("ebooklib and beautifulsoup4 are required to parse EPUB.") from exc

        book = epub.read_epub(str(file_ref.path))
        texts: List[str] = []

        for item in book.get_items():
            if item.get_type() == epub.ITEM_DOCUMENT:
                content = item.get_body_content().decode("utf-8", errors="ignore")
                soup = BeautifulSoup(content, "html.parser")
                page_text = soup.get_text(separator="\n")
                texts.append(page_text)

        full_text = "\n\n".join(texts)
        chunks = simple_text_chunker(full_text)
        meta = {
            "file_id": file_ref.id,
            "file_type": file_ref.file_type.value,
            "parser": "EPUBParser",
        }
        return DocumentText(text=full_text, chunks=chunks, metadata=meta)


# ---------- Image OCR ----------

class ImageOCRParser(BaseDocumentParser):
    file_types = [FileType.IMAGE_PNG, FileType.IMAGE_JPEG, FileType.IMAGE_WEBP]

    def parse(self, file_ref: FileRef) -> DocumentText:
        try:
            from PIL import Image  # type: ignore
            import pytesseract  # type: ignore
        except ImportError as exc:  # pragma: no cover
            raise InfraError(
                "pillow and pytesseract are required for ImageOCRParser."
            ) from exc

        img = Image.open(file_ref.path)
        text = pytesseract.image_to_string(img)
        chunks = simple_text_chunker(text)
        meta = {
            "file_id": file_ref.id,
            "file_type": file_ref.file_type.value,
            "parser": "ImageOCRParser",
        }
        return DocumentText(text=text, chunks=chunks, metadata=meta)


# ---------- dispatcher ----------

class DocumentParserDispatcher:
    """
    Central dispatcher that picks the right parser based on FileType.
    """

    def __init__(self) -> None:
        self._parsers: List[BaseDocumentParser] = [
            PlainTextParser(),
            MarkdownParser(),
            PDFParser(),
            DOCXParser(),
            PPTXParser(),
            HTMLParser(),
            EPUBParser(),
            ImageOCRParser(),
        ]

    def register_parser(self, parser: BaseDocumentParser) -> None:
        self._parsers.append(parser)

    def parse(self, file_ref: FileRef) -> DocumentText:
        for parser in self._parsers:
            if parser.can_handle(file_ref):
                logger.debug(
                    "DocumentParserDispatcher.parse",
                    extra={
                        "file_id": file_ref.id,
                        "file_type": file_ref.file_type.value,
                        "parser": type(parser).__name__,
                    },
                )
                return parser.parse(file_ref)
        raise InfraError(f"No parser registered for file type {file_ref.file_type.value}.")


# ---------- pickling helpers ----------

def pickle_document_text(doc: DocumentText, target_path: Path) -> None:
    with target_path.open("wb") as f:
        pickle.dump(doc, f)


def load_pickled_document_text(source_path: Path) -> DocumentText:
    with source_path.open("rb") as f:
        obj = pickle.load(f)
    if not isinstance(obj, DocumentText):
        raise InfraError("Loaded object is not a DocumentText.")
    return obj